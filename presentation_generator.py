from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime


def create_presentation(results, output_file):
    """Создаёт презентацию для руководства"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === СЛАЙД 1: ТИТУЛЬНЫЙ ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой макет

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "АНАЛИЗ ДОГОВОРОВ"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(68, 114, 196)
    title_para.alignment = PP_ALIGN.CENTER

    # Подзаголовок
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Отчёт сформирован: {datetime.now().strftime('%d.%m.%Y')}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Статистика
    stats_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    stats_frame = stats_box.text_frame
    stats_frame.text = f"Обработано договоров: {len(results)}"
    stats_para = stats_frame.paragraphs[0]
    stats_para.font.size = Pt(28)
    stats_para.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 2: СВОДКА ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "СВОДНАЯ ИНФОРМАЦИЯ"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(68, 114, 196)

    # Статистика
    total_contracts = len(results)
    high_risk = sum(1 for r in results if 'ВЫШЕ НОРМЫ' in r['peni_risk'])
    normal_risk = sum(1 for r in results if 'норма' in r['peni_risk'])
    no_data = total_contracts - high_risk - normal_risk

    stats_text = f"""
Всего договоров: {total_contracts}

🔴 Высокий риск (пени > 0.1%): {high_risk}
🟢 Норма (пени ≤ 0.1%): {normal_risk}
⚪ Нет данных: {no_data}
    """

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = stats_text
    content_frame.word_wrap = True

    for para in content_frame.paragraphs:
        para.font.size = Pt(24)
        para.space_after = Pt(12)

    # === СЛАЙДЫ ПО КАЖДОМУ ДОГОВОРУ ===
    for idx, data in enumerate(results, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Заголовок с номером
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = f"ДОГОВОР {idx}: {data['filename'][:50]}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(68, 114, 196)

        # Информация о договоре
        info_text = f"""
📄 Номер: {data['number']}
📅 Дата: {data['date']}
💰 Сумма: {data['amount']}
🏢 ИНН: {len(data['inn_list'])} шт.
📊 Пени: {data['peni']} — {data['peni_risk']}
        """

        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(2.5))
        info_frame = info_box.text_frame
        info_frame.text = info_text
        info_frame.word_wrap = True

        for para in info_frame.paragraphs:
            para.font.size = Pt(18)
            para.space_after = Pt(8)

        # ИИ-анализ
        ai_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(3))
        ai_frame = ai_box.text_frame
        ai_frame.text = "🤖 АНАЛИЗ РИСКОВ (ИИ):\n\n" + data.get('ai_analysis', 'Нет данных')
        ai_frame.word_wrap = True

        ai_title = ai_frame.paragraphs[0]
        ai_title.font.size = Pt(20)
        ai_title.font.bold = True
        ai_title.font.color.rgb = RGBColor(68, 114, 196)

        for para in ai_frame.paragraphs[1:]:
            para.font.size = Pt(16)

    # === ФИНАЛЬНЫЙ СЛАЙД ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    final_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    final_frame = final_box.text_frame
    final_frame.text = "СПАСИБО ЗА ВНИМАНИЕ"
    final_para = final_frame.paragraphs[0]
    final_para.font.size = Pt(48)
    final_para.font.bold = True
    final_para.font.color.rgb = RGBColor(68, 114, 196)
    final_para.alignment = PP_ALIGN.CENTER

    # Сохраняем
    prs.save(output_file)
    return output_file